#!/usr/bin/env python3
"""Extract text and images from .pptx files.

Uses python-pptx to iterate slides, extract text from shapes,
and save embedded images to an assets directory.

Usage: python _pptx_extract.py <pptx_path> <assets_dir>
Output: JSON line with {markdown, slides, images} on stdout.
"""
import json
import os
import sys
from pathlib import Path

def extract_pptx(pptx_path: str, assets_dir: str) -> dict:
    try:
        from pptx import Presentation
    except ImportError:
        return {"error": "python-pptx not installed", "markdown": "", "slides": 0, "images": 0}

    prs = Presentation(pptx_path)
    sections: list[str] = []
    image_count = 0
    os.makedirs(assets_dir, exist_ok=True)

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_title = ""
        slide_texts: list[str] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    if not slide_title and len(text) < 120:
                        slide_title = text
                    else:
                        slide_texts.append(text)

            if hasattr(shape, "image"):
                try:
                    img = shape.image
                    ext_map = {
                        "image/png": ".png",
                        "image/jpeg": ".jpg",
                        "image/gif": ".gif",
                        "image/bmp": ".bmp",
                        "image/svg+xml": ".svg",
                        "image/webp": ".webp",
                    }
                    ext = ext_map.get(img.content_type, ".png")
                    img_name = f"slide{slide_num:03d}-img{image_count + 1:02d}{ext}"
                    img_path = os.path.join(assets_dir, img_name)
                    with open(img_path, "wb") as f:
                        f.write(img.blob)
                    image_count += 1
                    slide_texts.append(f"![Slide {slide_num} Image](assets/{img_name})")
                except Exception:
                    pass

        heading = slide_title or f"Slide {slide_num}"
        sections.append(f"## {heading}\n")
        if slide_texts:
            sections.append("\n".join(slide_texts))
        sections.append("")

    markdown = "\n".join(sections).strip()
    return {
        "markdown": markdown,
        "slides": len(prs.slides),
        "images": image_count,
    }


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print(json.dumps({"error": "Usage: _pptx_extract.py <pptx_path> <assets_dir>"}))
        sys.exit(1)

    result = extract_pptx(sys.argv[1], sys.argv[2])
    print(json.dumps(result, ensure_ascii=False))
